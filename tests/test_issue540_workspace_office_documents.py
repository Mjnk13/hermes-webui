from __future__ import annotations

import io

import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from api.office_documents import (
    CLAIMED_OFFICE_EXTENSIONS,
    is_claimed_office_path,
    preview_office_document,
    save_office_document,
)


def _simple_docx_bytes(*paragraphs: str) -> bytes:
    document = DocxDocument()
    for paragraph in paragraphs or ("",):
        document.add_paragraph(paragraph)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _rich_docx_bytes() -> bytes:
    document = DocxDocument()
    document.add_paragraph("Lead paragraph")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "left"
    table.cell(0, 1).text = "right"
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _simple_xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Summary"
    sheet["A1"] = "alpha"
    sheet["B1"] = "beta"
    sheet["A2"] = "gamma"
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _simple_pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text = "Office preview"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_office_registry_claims_docx_xlsx_pptx():
    assert CLAIMED_OFFICE_EXTENSIONS == {".docx", ".xlsx", ".pptx"}
    assert is_claimed_office_path("report.docx")
    assert is_claimed_office_path("budget.xlsx")
    assert is_claimed_office_path("deck.pptx")

    docx_preview = preview_office_document("report.docx", _simple_docx_bytes("one", "two"))
    xlsx_preview = preview_office_document("budget.xlsx", _simple_xlsx_bytes())
    pptx_preview = preview_office_document("deck.pptx", _simple_pptx_bytes())

    assert docx_preview["preview_kind"] == "office"
    assert xlsx_preview["preview_kind"] == "office"
    assert pptx_preview["preview_kind"] == "office"
    assert docx_preview["office_format"] == "docx"
    assert xlsx_preview["office_format"] == "xlsx"
    assert pptx_preview["office_format"] == "pptx"
    assert docx_preview["render_mode"] == "code"
    assert xlsx_preview["render_mode"] == "code"
    assert pptx_preview["render_mode"] == "code"
    assert docx_preview["editable"] is True
    assert xlsx_preview["editable"] is False
    assert pptx_preview["editable"] is False


def test_docx_paragraph_projection_round_trips_simple_documents():
    original_bytes = _simple_docx_bytes("alpha", "beta")
    preview = preview_office_document("story.docx", original_bytes)

    assert preview["editable"] is True
    assert preview["content"] == "alpha\n\nbeta"

    saved_preview, saved_bytes = save_office_document("story.docx", original_bytes, "alpha\nbeta\ngamma")
    round_trip = DocxDocument(io.BytesIO(saved_bytes))

    assert saved_preview["editable"] is True
    assert saved_preview["content"] == "alpha\nbeta\ngamma"
    assert [paragraph.text for paragraph in round_trip.paragraphs] == ["alpha", "beta", "gamma"]


def test_xlsx_stays_preview_only():
    path = "budget.xlsx"
    raw = _simple_xlsx_bytes()

    preview = preview_office_document(path, raw)

    assert preview["preview_kind"] == "office"
    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")

    with pytest.raises(ValueError):
        save_office_document(path, raw, "edited text")


def test_pptx_stays_preview_only():
    path = "deck.pptx"
    raw = _simple_pptx_bytes()

    preview = preview_office_document(path, raw)

    assert preview["preview_kind"] == "office"
    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")

    with pytest.raises(ValueError):
        save_office_document(path, raw, "edited text")


def test_rich_docx_stays_preview_only():
    path = "rich.docx"
    raw = _rich_docx_bytes()

    preview = preview_office_document(path, raw)

    assert preview["preview_kind"] == "office"
    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")

    with pytest.raises(ValueError):
        save_office_document(path, raw, "edited text")
